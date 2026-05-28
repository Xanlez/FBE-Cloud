"""Собрать fbe-cloud-presentation.pptx в дизайне сайта FBE cloud."""
import re
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = Path(__file__).resolve().parent / "fbe-cloud-presentation.pptx"
LOGO = ROOT / "logo_FBE.png"

# Палитра из static/css/app.css
BG = RGBColor(0x0F, 0x14, 0x19)
SURFACE = RGBColor(0x1A, 0x23, 0x32)
BORDER = RGBColor(0x2D, 0x3A, 0x4D)
TEXT = RGBColor(0xE8, 0xED, 0xF4)
MUTED = RGBColor(0x8B, 0x9C, 0xB3)
ACCENT = RGBColor(0x5B, 0x9F, 0xD4)
GRAD_TOP = RGBColor(0x1E, 0x3A, 0x5F)

# Стандарт PowerPoint «Широкоэкранный 16:9» (33.867 × 19.05 см)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
FONT = "Segoe UI"


def _apply_slide_size_16x9(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def _fix_presentation_xml_16x9(path: Path) -> None:
    """python-pptx оставляет type=screen4x3 — правим на screen16x9 для PowerPoint."""
    tmp = path.with_suffix(".tmp.pptx")
    sld_sz = '<p:sldSz cx="12192000" cy="6858000" type="screen16x9"/>'
    with zipfile.ZipFile(path, "r") as zin:
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if info.filename == "ppt/presentation.xml":
                    xml = data.decode("utf-8")
                    xml = re.sub(r"<p:sldSz[^>]*/>", sld_sz, xml, count=1)
                    data = xml.encode("utf-8")
                zout.writestr(info, data)
    shutil.move(tmp, path)


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _line(shape, color, width_pt=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.8)
    )
    _fill(band, GRAD_TOP)
    band.line.fill.background()
    band.fill.transparency = 0.35


def _header(slide, n: int, total: int):
    bar_h = Inches(0.72)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, bar_h)
    _fill(bar, SURFACE)
    _line(bar, BORDER)
    bar.line.fill.background()

    if LOGO.exists():
        slide.shapes.add_picture(
            str(LOGO), Inches(0.35), Inches(0.1), height=Inches(0.52)
        )

    brand = slide.shapes.add_textbox(Inches(1.05), Inches(0.18), Inches(2), Inches(0.4))
    p = brand.text_frame.paragraphs[0]
    p.text = "FBE cloud"
    p.font.name = FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT

    counter = slide.shapes.add_textbox(Inches(11.5), Inches(0.22), Inches(1.5), Inches(0.35))
    p = counter.text_frame.paragraphs[0]
    p.text = f"{n} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED


def _footer(slide):
    foot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.88), SLIDE_W, Inches(0.62)
    )
    _fill(foot, SURFACE)
    _line(foot, BORDER)

    left = slide.shapes.add_textbox(Inches(0.4), Inches(7.02), Inches(5), Inches(0.35))
    p = left.text_frame.paragraphs[0]
    p.text = "FBE cloud v1.72pr · presentation version"
    p.font.name = FONT
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED

    right = slide.shapes.add_textbox(Inches(9.5), Inches(7.02), Inches(3.4), Inches(0.35))
    p = right.text_frame.paragraphs[0]
    p.text = "by Xanlez, 2026"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def _title(slide, text: str, top=Inches(1.15), size=32, center=False):
    box = slide.shapes.add_textbox(
        Inches(0.75), top, Inches(11.8), Inches(0.9)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TEXT
    if center:
        p.alignment = PP_ALIGN.CENTER


def _body(slide, text: str, top, height=Inches(0.9), center=False, muted=True):
    box = slide.shapes.add_textbox(Inches(0.75), top, Inches(11.8), height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(14)
    p.font.color.rgb = MUTED if muted else TEXT
    if center:
        p.alignment = PP_ALIGN.CENTER


def _bullets(slide, items: list[str], top=Inches(1.95)):
    box = slide.shapes.add_textbox(Inches(0.9), top, Inches(11.5), Inches(2.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED
        p.space_after = Pt(6)


def _card(slide, left, top, width, height, title: str, body: str):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(rect, SURFACE)
    _line(rect, BORDER)
    rect.adjustments[0] = 0.08

    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), height - Inches(0.2)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    h = tf.paragraphs[0]
    h.text = title
    h.font.name = FONT
    h.font.size = Pt(12)
    h.font.bold = True
    h.font.color.rgb = ACCENT
    h.space_after = Pt(4)

    b = tf.add_paragraph()
    b.text = body
    b.font.name = FONT
    b.font.size = Pt(10)
    b.font.color.rgb = MUTED


def _tags(slide, tags: list[str], top=Inches(4.2)):
    x = Inches(0.75)
    for tag in tags:
        w = Inches(0.55 + len(tag) * 0.09)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, Inches(0.38))
        _fill(pill, RGBColor(0x14, 0x28, 0x3D))
        _line(pill, BORDER)
        pill.adjustments[0] = 0.5
        tb = slide.shapes.add_textbox(x, top + Inches(0.06), w, Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = tag
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        x += w + Inches(0.12)


def _slide_base(prs, n: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _header(slide, n, total)
    _footer(slide)
    return slide


def build():
    prs = Presentation()
    _apply_slide_size_16x9(prs)
    total = 8

    # 1 — титул
    s = _slide_base(prs, 1, total)
    if LOGO.exists():
        s.shapes.add_picture(
            str(LOGO), Inches(4.2), Inches(1.35), width=Inches(4.9)
        )
    _title(s, "Учебные материалы всегда под рукой", Inches(3.55), 30, center=True)
    _body(
        s,
        "FBE cloud — облако для конспектов, файлов курсов и совместной работы студентов.",
        Inches(4.45),
        Inches(0.7),
        center=True,
    )
    _body(s, "v1.72pr · presentation version · 2026", Inches(5.35), Inches(0.4), center=True)

    # 2 — о проекте
    s = _slide_base(prs, 2, total)
    _title(s, "О проекте")
    _body(s, "Веб-платформа для хранения и обмена учебными материалами в одном месте.", Inches(1.85))
    _card(
        s,
        Inches(0.75),
        Inches(2.55),
        Inches(11.8),
        Inches(1.55),
        "Задача",
        "Разрозненные файлы в мессенджерах и на дисках сложно искать и делить с одногруппниками. "
        "FBE cloud объединяет личное хранилище, общее облако и социальные функции.",
    )

    # 3 — разделы
    s = _slide_base(prs, 3, total)
    _title(s, "Основные разделы")
    cards = [
        ("Мои файлы", "Личное облако: загрузка, предпросмотр, фильтры и управление своими материалами."),
        ("Общее облако", "Общая библиотека файлов для всех пользователей с поиском и фильтрацией."),
        ("Друзья и чат", "Добавление друзей и обмен сообщениями внутри платформы."),
        ("События", "Экзамены, зачёты и мероприятия — открытые и закрытые, с датами и описанием."),
    ]
    positions = [
        (Inches(0.75), Inches(1.75)),
        (Inches(6.55), Inches(1.75)),
        (Inches(0.75), Inches(3.85)),
        (Inches(6.55), Inches(3.85)),
    ]
    for (left, top), (title, body) in zip(positions, cards):
        _card(s, left, top, Inches(5.7), Inches(1.85), title, body)

    # 4 — аккаунты
    s = _slide_base(prs, 4, total)
    _title(s, "Аккаунты и безопасность")
    _bullets(
        s,
        [
            "Регистрация и вход с сессиями",
            "Подтверждение e-mail",
            "Восстановление пароля",
            "Профиль с аватаром и сменой пароля",
            "Роли персонала и панель администратора",
        ],
    )
    _card(
        s,
        Inches(0.75),
        Inches(4.35),
        Inches(11.8),
        Inches(1.1),
        "Админка",
        "Статистика пользователей и управление для сотрудников (personnel).",
    )

    # 5 — файлы
    s = _slide_base(prs, 5, total)
    _title(s, "Работа с файлами")
    _bullets(
        s,
        [
            "Загрузка с лимитами на объём и количество",
            "Предпросмотр документов и изображений в браузере",
            "Фильтры по типу, дате и названию",
            "Интеграция с Google Drive (опционально)",
        ],
    )
    _tags(s, ["FastAPI", "SQLAlchemy", "Jinja2"])

    # 6 — дизайн
    s = _slide_base(prs, 6, total)
    _title(s, "Дизайн интерфейса")
    _body(s, "Тёмная тема в стиле этой презентации и сайта.", Inches(1.85))
    design_cards = [
        ("Палитра", "Фон #0f1419, поверхности #1a2332, акцент #5b9fd4, текст #e8edf4."),
        ("Компоненты", "Карточки, кнопки primary/ghost, боковое меню, explorer для файлов."),
        ("Шрифт", "Segoe UI / system-ui — как на сайте."),
    ]
    for i, (title, body) in enumerate(design_cards):
        _card(s, Inches(0.75 + i * 4.05), Inches(2.45), Inches(3.85), Inches(2.2), title, body)

    # 7 — технологии
    s = _slide_base(prs, 7, total)
    _title(s, "Технологии")
    _card(
        s,
        Inches(0.75),
        Inches(1.75),
        Inches(5.7),
        Inches(3.5),
        "Backend",
        "• Python, FastAPI, Uvicorn\n• SQLAlchemy, SQLite\n• Session middleware, proxy headers",
    )
    _card(
        s,
        Inches(6.55),
        Inches(1.75),
        Inches(5.7),
        Inches(3.5),
        "Frontend",
        "• Шаблоны Jinja2\n• CSS custom properties\n• Адаптивная вёрстка",
    )

    # 8 — финал
    s = _slide_base(prs, 8, total)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(5.4), Inches(1.5), height=Inches(1.4))
    _title(s, "Спасибо за внимание", Inches(3.2), 32, center=True)
    _body(
        s,
        "FBE cloud — учебные материалы всегда под рукой.",
        Inches(4.15),
        Inches(0.6),
        center=True,
    )
    _body(s, "by Xanlez, 2026", Inches(5.0), Inches(0.4), center=True)

    built = OUT.with_name("_build_fbe-cloud-presentation.pptx")
    prs.save(built)
    _fix_presentation_xml_16x9(built)
    try:
        if OUT.exists():
            OUT.unlink()
        shutil.move(built, OUT)
        target = OUT
    except OSError:
        target = built
        print(f"Close {OUT.name} in PowerPoint and rerun the script.")
    print(f"Saved (16:9): {target}")


if __name__ == "__main__":
    build()
